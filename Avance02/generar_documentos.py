"""
Generador automático de Word (.docx) y PowerPoint (.pptx)
para el Avance 02 del proyecto GymMax.

Ejecutar con: python generar_documentos.py
"""

from docx import Document
from docx.shared import Pt, Cm, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_ALIGN_VERTICAL
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.enum.text import PP_ALIGN

import os

# ===== CONFIGURACIÓN =====
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
OUT_DOCX = os.path.join(BASE_DIR, "GymMax_Avance02.docx")
OUT_PPTX = os.path.join(BASE_DIR, "GymMax_Avance02.pptx")

ORANGE = RGBColor(0xFF, 0x6B, 0x00)
DARK   = RGBColor(0x1A, 0x1A, 0x1A)
GRAY   = RGBColor(0x55, 0x55, 0x55)

# ============================================================
#                 GENERACIÓN DEL WORD (.docx)
# ============================================================

def set_cell_bg(cell, hex_color):
    """Pinta el fondo de una celda."""
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), hex_color)
    tc_pr.append(shd)


def add_heading(doc, text, level=1, color=None, center=False):
    h = doc.add_heading(text, level=level)
    if center:
        h.alignment = WD_ALIGN_PARAGRAPH.CENTER
    if color:
        for run in h.runs:
            run.font.color.rgb = color
    return h


def add_para(doc, text, bold=False, italic=False, size=11, color=None, align=None):
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.bold = bold
    run.italic = italic
    run.font.size = Pt(size)
    if color:
        run.font.color.rgb = color
    if align is not None:
        p.alignment = align
    return p


def add_bullet(doc, text, bold_prefix=None):
    p = doc.add_paragraph(style='List Bullet')
    if bold_prefix:
        run = p.add_run(bold_prefix)
        run.bold = True
        p.add_run(" " + text)
    else:
        p.add_run(text)
    return p


def add_table_from_data(doc, headers, rows, header_bg="FF6B00", header_white=True):
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.style = "Light Grid Accent 1"
    # Headers
    hdr_cells = table.rows[0].cells
    for i, h in enumerate(headers):
        hdr_cells[i].text = h
        set_cell_bg(hdr_cells[i], header_bg)
        for p in hdr_cells[i].paragraphs:
            for run in p.runs:
                run.bold = True
                if header_white:
                    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                run.font.size = Pt(10)
    # Data
    for r_idx, row in enumerate(rows, start=1):
        cells = table.rows[r_idx].cells
        for c_idx, val in enumerate(row):
            cells[c_idx].text = str(val)
            for p in cells[c_idx].paragraphs:
                for run in p.runs:
                    run.font.size = Pt(10)
    return table


def add_page_break(doc):
    doc.add_page_break()


# ----- CARÁTULA -----
def build_caratula(doc):
    for _ in range(3):
        doc.add_paragraph()
    add_para(doc, "UNIVERSIDAD TECNOLÓGICA DEL PERÚ", bold=True, size=16,
             align=WD_ALIGN_PARAGRAPH.CENTER, color=DARK)
    add_para(doc, "Facultad de Ingeniería", size=12, align=WD_ALIGN_PARAGRAPH.CENTER)
    add_para(doc, "Carrera de Ingeniería de Sistemas e Informática", size=12,
             align=WD_ALIGN_PARAGRAPH.CENTER)
    for _ in range(3):
        doc.add_paragraph()
    add_para(doc, "INFORME DE AVANCE 02", bold=True, size=22,
             align=WD_ALIGN_PARAGRAPH.CENTER, color=ORANGE)
    add_para(doc, "Diseño y Avance de Funcionalidad", italic=True, size=14,
             align=WD_ALIGN_PARAGRAPH.CENTER, color=GRAY)
    doc.add_paragraph()
    add_para(doc, "GymMax", bold=True, size=28,
             align=WD_ALIGN_PARAGRAPH.CENTER, color=ORANGE)
    add_para(doc, "Sistema Integral de Gestión de Gimnasios", italic=True, size=12,
             align=WD_ALIGN_PARAGRAPH.CENTER)
    for _ in range(2):
        doc.add_paragraph()
    add_para(doc, "Curso: Desarrollo Web Integrado", bold=True, size=12,
             align=WD_ALIGN_PARAGRAPH.CENTER)
    add_para(doc, "Sección: 27672", size=11, align=WD_ALIGN_PARAGRAPH.CENTER)
    add_para(doc, "Docente: Juan Manuel Rodríguez del Águila", size=11,
             align=WD_ALIGN_PARAGRAPH.CENTER)
    doc.add_paragraph()
    add_para(doc, "Integrantes (Grupo 3):", bold=True, size=12,
             align=WD_ALIGN_PARAGRAPH.CENTER)
    integrantes = [
        "Choque Anchante, Niurka Yasbeth",
        "Diaz Culqui, Nehemias",
        "Purizaca Ipanaque, Dennys Marlon",
        "Soria Chavez, Ian",
        "Torre Escobar, Oliver",
        "Valladolid Llenque, Alexander",
    ]
    for nombre in integrantes:
        add_para(doc, nombre, size=11, align=WD_ALIGN_PARAGRAPH.CENTER)
    for _ in range(3):
        doc.add_paragraph()
    add_para(doc, "Lima, Mayo 2026", size=11, align=WD_ALIGN_PARAGRAPH.CENTER, color=GRAY)
    add_page_break(doc)


# ----- ÍNDICE (placeholder) -----
def build_indice(doc):
    add_heading(doc, "Índice", level=1, color=ORANGE)
    add_para(doc, "Capítulo I: Introducción ........................................ 3", size=11)
    add_para(doc, "    1.1. Introducción general", size=11)
    add_para(doc, "    1.2. El producto", size=11)
    add_para(doc, "    1.3. Antecedentes y problemática", size=11)
    add_para(doc, "    1.4. Objetivos", size=11)
    add_para(doc, "    1.5. Alcance", size=11)
    add_para(doc, "    1.6. Características del producto", size=11)
    doc.add_paragraph()
    add_para(doc, "Capítulo II: Arquitectura", size=11, bold=True)
    add_para(doc, "    2.1. Product Backlog", size=11)
    add_para(doc, "    2.2. User Stories", size=11)
    add_para(doc, "    2.3. Diagrama de flujo de navegación", size=11)
    add_para(doc, "    2.4. Wireframes (Mobile y Desktop)", size=11)
    add_para(doc, "    2.5. Modelo de Base de Datos (Lógico y Físico)", size=11)
    add_para(doc, "    2.6. Tecnologías aplicadas", size=11)
    add_para(doc, "    2.7. Enlaces a artefactos", size=11)
    doc.add_paragraph()
    add_para(doc, "Capítulo III: Backend", size=11, bold=True)
    add_para(doc, "    3.1. Diagrama de Clases UML", size=11)
    add_para(doc, "    3.2. Diccionario de Clases", size=11)
    add_page_break(doc)


# ----- CAPÍTULO I -----
def build_cap1(doc):
    add_heading(doc, "Capítulo I: Introducción", level=1, color=ORANGE)

    add_heading(doc, "1.1. Introducción general", level=2)
    add_para(doc,
        "En el contexto actual del mercado peruano, el sector de gimnasios y centros de "
        "acondicionamiento físico ha experimentado un crecimiento sostenido en los últimos "
        "años, especialmente tras la pandemia, donde el cuidado de la salud se convirtió "
        "en una prioridad para gran parte de la población. Sin embargo, este crecimiento "
        "ha evidenciado las limitaciones de los métodos tradicionales de gestión: registros "
        "manuales en cuadernos, control de membresías mediante hojas de cálculo y atención "
        "en recepción que genera largas colas en horas pico.")
    add_para(doc,
        "Frente a esta realidad, el presente trabajo propone el desarrollo de GymMax, una "
        "aplicación web integral construida sobre la arquitectura Java Enterprise Edition "
        "(Jakarta EE) que permitirá a las cadenas de gimnasios automatizar y centralizar "
        "la gestión de sus sedes, socios, membresías, clases, reservas y reportes financieros. "
        "La solución se diseña con un enfoque Mobile First utilizando Bootstrap 5.3.")
    add_para(doc,
        "El documento detalla el avance del diseño y la arquitectura de la solución, "
        "abarcando desde la identificación de los actores y casos de uso hasta el modelado "
        "completo de la base de datos, el diagrama de clases UML del backend y el "
        "diccionario detallado de cada componente del sistema.")

    add_heading(doc, "1.2. El producto", level=2)
    add_heading(doc, "1.2.1. Nombre del producto", level=3)
    add_para(doc,
        "El producto se denomina GymMax — Sistema Integral de Gestión de Gimnasios. "
        "El nombre combina los términos Gym (gimnasio) y Max (máximo), transmitiendo "
        "la propuesta de valor de 'llevar al máximo la gestión del gimnasio'.")
    add_heading(doc, "1.2.2. Descripción del producto", level=3)
    add_para(doc,
        "GymMax es una plataforma web centralizada que digitaliza el ciclo de vida completo "
        "del socio de un gimnasio: desde el registro inicial y la contratación de membresías "
        "digitales, hasta la reserva online de clases en cualquiera de las sedes de la cadena. "
        "El sistema incorpora también un panel administrativo robusto con KPIs en tiempo real, "
        "gestión CRUD de socios y sedes, generación de reportes financieros y exportación a "
        "Excel y PDF.")

    add_heading(doc, "1.3. Antecedentes y problemática", level=2)
    add_heading(doc, "1.3.1. Antecedentes del sector", level=3)
    add_para(doc, "SmartGym Perú (2022) — Lima.", bold=True)
    add_para(doc,
        "Desarrollado por estudiantes de la Universidad Nacional Mayor de San Marcos (UNMSM), "
        "este proyecto consistió en un sistema web de gestión de membresías para una cadena con "
        "tres sedes en Lima Metropolitana. Implementado en PHP y MySQL, logró reducir en un 58% "
        "el tiempo de registro de nuevos socios y disminuir la morosidad en pagos mensuales en "
        "un 35%.")
    add_para(doc, "GymSys Trujillo (2020) — La Libertad.", bold=True)
    add_para(doc,
        "Desarrollado por la Universidad César Vallejo (UCV) como tesis de pregrado, GymSys fue "
        "inicialmente un sistema de escritorio convertido posteriormente a web para gestionar "
        "reservas de clases grupales (spinning, yoga, crossfit). Incrementó la ocupación de "
        "horarios pico en un 40%.")
    add_para(doc, "BodyTrack (2023) — Callao.", bold=True)
    add_para(doc,
        "Proyecto desarrollado por egresados de la Universidad Tecnológica del Perú (UTP) para "
        "un gimnasio mediano del Callao. BodyTrack implementó control de pagos recurrentes y "
        "reportes de ingresos, logrando una disminución del 50% en la cartera morosa y un "
        "incremento del 25% en la retención de socios.")

    add_heading(doc, "1.3.2. Problemática identificada", level=3)
    add_para(doc, "Los gimnasios peruanos continúan enfrentando los siguientes problemas:")
    add_bullet(doc, "Sin un sistema automatizado, las membresías vencen sin que el gimnasio detecte oportunamente la oportunidad de renovación.",
               "Pérdida de ingresos por falta de control de vencimientos.")
    add_bullet(doc, "El personal de recepción dedica horas diarias a registrar pagos manualmente.",
               "Sobrecarga administrativa.")
    add_bullet(doc, "En horarios pico, los socios deben esperar para confirmar asistencia o reservar clases.",
               "Procesos lentos y colas en recepción.")
    add_bullet(doc, "La imposibilidad de revisar membresía o reservar clases desde el celular obliga al socio a acudir físicamente al gimnasio.",
               "Mala experiencia del socio.")
    add_bullet(doc, "Los pocos sistemas comerciales disponibles suelen presentar lentitud y caídas en horarios de alta demanda.",
               "Sistemas existentes con bajo rendimiento.")

    add_heading(doc, "1.4. Objetivos", level=2)
    add_heading(doc, "1.4.1. Objetivo general", level=3)
    add_para(doc,
        "Desarrollar una aplicación web integral bajo la arquitectura Java Enterprise Edition "
        "(Jakarta EE) que centralice la visualización de sedes, el registro de socios, la venta "
        "de membresías digitales y la gestión de reservas de clases, garantizando una experiencia "
        "Mobile First mediante Bootstrap 5.3 y persistencia segura en MySQL.")
    add_heading(doc, "1.4.2. Objetivos específicos", level=3)
    objs = [
        "Implementar un módulo de consulta de sedes que permita visualizar ubicación, horarios, servicios y cupos por clase.",
        "Desarrollar un formulario web responsive de registro de socios con validación en tiempo real y almacenamiento seguro en MySQL.",
        "Crear un sistema de gestión de planes (Básico, Premium, Anual) con cálculo automático de vencimiento y generación de comprobante.",
        "Diseñar una interfaz Mobile First con Bootstrap 5.3 accesible desde PC, tablet y smartphone.",
        "Implementar un panel administrativo con KPIs en tiempo real, CRUD de socios y sedes, y exportación de reportes a Excel y PDF.",
        "Aplicar la arquitectura por capas (MVC + DAO + Facade) para garantizar mantenibilidad y escalabilidad.",
    ]
    for i, o in enumerate(objs, 1):
        add_para(doc, f"{i}. {o}", size=11)

    add_heading(doc, "1.5. Alcance", level=2)
    add_para(doc,
        "El alcance de GymMax contempla el desarrollo de una plataforma web centralizada que "
        "permite la gestión completa del ciclo de vida del socio.")
    add_para(doc, "Funcionalidades incluidas:", bold=True)
    add_bullet(doc, "Registro y autenticación de socios y administradores con roles diferenciados.")
    add_bullet(doc, "Gestión multi-sede con datos independientes por local.")
    add_bullet(doc, "Catálogo de planes (Básico, Premium, Anual) con renovación autónoma.")
    add_bullet(doc, "Reserva y cancelación de clases con control de cupos en tiempo real.")
    add_bullet(doc, "Procesamiento de pagos digitales (Yape, Plin, Tarjeta) con comprobantes.")
    add_bullet(doc, "Panel administrativo con dashboard de KPIs, CRUD y reportes exportables.")
    add_bullet(doc, "Interfaz responsive Mobile First.")
    add_para(doc, "Limitaciones del alcance:", bold=True)
    add_bullet(doc, "No incluye pagos en efectivo (solo digitales).")
    add_bullet(doc, "No incluye aplicación móvil nativa.")
    add_bullet(doc, "No incluye integración biométrica.")
    add_bullet(doc, "No incluye marketing por correo masivo.")
    add_para(doc, "Público objetivo:", bold=True)
    add_para(doc,
        "Cadenas de gimnasios peruanos (2-10 sedes) que buscan migrar de administración manual "
        "a infraestructura digital basada en Jakarta EE.")

    add_heading(doc, "1.6. Características del producto", level=2)
    add_heading(doc, "1.6.1. Características funcionales", level=3)
    caracs_func = [
        ("Gestión multi-sede centralizada.", "Administración unificada de múltiples sedes bajo una sola base de datos centralizada."),
        ("Sistema de roles diferenciados.", "Dos perfiles: Socio (cliente) y Administrador (personal), cada uno con flujo y permisos propios."),
        ("Registro y autenticación segura.", "Formulario responsive con validación en tiempo real y almacenamiento seguro de contraseñas."),
        ("Gestión de membresías digitales.", "Soporta planes Básico, Premium y Anual con cálculo automático de vencimiento."),
        ("Reserva de clases online.", "Visualización de cupos por clase, reserva y cancelación según políticas del gimnasio."),
        ("Control de asistencia y Kardex.", "Registro automatizado de ingresos y salidas consultable por socio y admin."),
        ("Procesamiento de pagos digitales.", "Validación de transacciones por Yape, Plin o Tarjeta con generación de comprobantes."),
        ("Dashboard administrativo con KPIs.", "Panel con socios activos, ingresos, vencimientos, reservas, gráficos y actividad reciente."),
        ("Reportes y exportación.", "Reportes filtrables por sede, fecha y tipo, exportables a Excel y PDF."),
        ("CRUD completo de entidades maestras.", "Crear/Leer/Actualizar/Eliminar socios, sedes, planes, clases y horarios."),
    ]
    for titulo, desc in caracs_func:
        add_bullet(doc, desc, titulo)

    add_heading(doc, "1.6.2. Características técnicas", level=3)
    caracs_tec = [
        ("Diseño Mobile First Responsive.", "Bootstrap 5.3 adaptable a móvil, tablet y escritorio."),
        ("Arquitectura por capas (MVC + DAO + Facade).", "Separación de responsabilidades y mantenibilidad."),
        ("Persistencia en MySQL.", "10 tablas normalizadas hasta 3FN con integridad referencial."),
        ("Despliegue en Apache Tomcat 10.", "Soporte para Jakarta EE 10."),
        ("Escalabilidad horizontal.", "Permite agregar sedes sin modificar la arquitectura."),
        ("Validación frontend y backend.", "Doble capa: HTML5+JS en cliente, Servlets en servidor."),
        ("Seguridad de sesión.", "Timeout configurable y control de acceso por rol."),
        ("Internacionalización de moneda.", "Soles peruanos (S/) con formato decimal de 2 posiciones."),
    ]
    for titulo, desc in caracs_tec:
        add_bullet(doc, desc, titulo)

    add_page_break(doc)


# ----- CAPÍTULO II -----
def build_cap2(doc):
    add_heading(doc, "Capítulo II: Arquitectura", level=1, color=ORANGE)

    # 2.1 Product Backlog
    add_heading(doc, "2.1. Product Backlog", level=2)
    add_para(doc,
        "El Product Backlog reúne todos los requerimientos funcionales y no funcionales "
        "identificados para el sistema GymMax, priorizados según valor de negocio y "
        "complejidad técnica. La estimación se realiza en Story Points (escala de Fibonacci).")

    add_heading(doc, "2.1.1. Requerimientos funcionales", level=3)
    rf_rows = [
        ("RF-01", "Registrar nuevo socio (formulario responsive)", "Autenticación", "Alta", "5", "1"),
        ("RF-02", "Iniciar sesión con correo y contraseña", "Autenticación", "Alta", "3", "1"),
        ("RF-03", "Recuperar contraseña vía correo", "Autenticación", "Media", "5", "3"),
        ("RF-04", "Validar campos del formulario en tiempo real", "Autenticación", "Alta", "3", "1"),
        ("RF-05", "Visualizar dashboard personal del socio", "Socio", "Alta", "5", "1"),
        ("RF-06", "Consultar listado de sedes disponibles", "Socio", "Alta", "3", "1"),
        ("RF-07", "Filtrar sedes por distrito", "Socio", "Media", "2", "2"),
        ("RF-08", "Ver detalle de servicios y cupos de una sede", "Socio", "Alta", "3", "2"),
        ("RF-09", "Reservar una clase en sede y horario", "Reservas", "Alta", "8", "2"),
        ("RF-10", "Visualizar listado de mis reservas", "Reservas", "Alta", "5", "2"),
        ("RF-11", "Cancelar una reserva confirmada", "Reservas", "Media", "3", "2"),
        ("RF-12", "Contratar un plan de membresía", "Membresías", "Alta", "8", "3"),
        ("RF-13", "Renovar membresía vencida o por vencer", "Membresías", "Alta", "5", "3"),
        ("RF-14", "Calcular vencimiento automático de membresía", "Membresías", "Alta", "2", "3"),
        ("RF-15", "Generar comprobante digital del pago", "Pagos", "Alta", "5", "3"),
        ("RF-16", "Procesar pagos por Yape, Plin o Tarjeta", "Pagos", "Alta", "8", "3"),
        ("RF-17", "Visualizar historial de pagos del socio", "Pagos", "Media", "3", "3"),
        ("RF-18", "Iniciar sesión como administrador (rol)", "Admin", "Alta", "3", "1"),
        ("RF-19", "Dashboard administrativo con KPIs", "Admin", "Alta", "8", "2"),
        ("RF-20", "CRUD de socios", "Admin", "Alta", "8", "2"),
        ("RF-21", "Filtrar socios por sede, plan y estado", "Admin", "Media", "3", "2"),
        ("RF-22", "CRUD de sedes", "Admin", "Alta", "5", "3"),
        ("RF-23", "Gestionar clases y horarios", "Admin", "Media", "5", "4"),
        ("RF-24", "Generar reportes de ingresos por sede", "Reportes", "Alta", "8", "4"),
        ("RF-25", "Exportar reportes a Excel y PDF", "Reportes", "Media", "5", "4"),
        ("RF-26", "Registrar asistencia del socio", "Asistencias", "Media", "5", "4"),
        ("RF-27", "Notificar vencimiento próximo al socio", "Notificaciones", "Baja", "8", "5"),
        ("RF-28", "Listar membresías por vencer al admin", "Admin", "Media", "3", "4"),
    ]
    add_table_from_data(doc,
        ["ID", "Requerimiento", "Módulo", "Prioridad", "SP", "Sprint"],
        rf_rows)

    add_heading(doc, "2.1.2. Requerimientos no funcionales", level=3)
    rnf_rows = [
        ("RNF-01", "Responsive Mobile/Tablet/Desktop con Bootstrap 5.3", "Usabilidad", "Alta"),
        ("RNF-02", "Tiempo de respuesta máximo 2s para CRUD", "Rendimiento", "Alta"),
        ("RNF-03", "Disponibilidad 99% en horarios del gimnasio", "Disponibilidad", "Alta"),
        ("RNF-04", "Contraseñas encriptadas (BCrypt)", "Seguridad", "Alta"),
        ("RNF-05", "Sesión HTTP con timeout 30 min", "Seguridad", "Alta"),
        ("RNF-06", "Soportar 200 usuarios concurrentes", "Escalabilidad", "Media"),
        ("RNF-07", "Compatibilidad navegadores modernos", "Compatibilidad", "Alta"),
        ("RNF-08", "Moneda en soles peruanos (2 decimales)", "Localización", "Media"),
        ("RNF-09", "Interfaz en español", "Localización", "Alta"),
        ("RNF-10", "Código versionado en Git", "Mantenibilidad", "Media"),
        ("RNF-11", "Arquitectura MVC + DAO + Facade", "Mantenibilidad", "Alta"),
        ("RNF-12", "Validación frontend + backend", "Seguridad", "Alta"),
    ]
    add_table_from_data(doc,
        ["ID", "Requerimiento", "Categoría", "Prioridad"],
        rnf_rows)

    add_heading(doc, "2.1.3. Distribución por Sprints", level=3)
    sprint_rows = [
        ("Sprint 1", "Autenticación + base del socio", "RF-01, 02, 04, 05, 06, 18", "22 SP"),
        ("Sprint 2", "Sedes, clases y reservas", "RF-07, 08, 09, 10, 11, 19, 20, 21", "37 SP"),
        ("Sprint 3", "Membresías y pagos", "RF-03, 12, 13, 14, 15, 16, 17, 22", "41 SP"),
        ("Sprint 4", "Reportes y asistencias", "RF-23, 24, 25, 26, 28", "26 SP"),
        ("Sprint 5", "Notificaciones y mejoras", "RF-27", "8 SP"),
    ]
    add_table_from_data(doc, ["Sprint", "Foco", "RFs", "Total SP"], sprint_rows)
    add_para(doc, "Total general: 134 Story Points distribuidos en 5 sprints.", italic=True)
    add_page_break(doc)

    # 2.2 User Stories
    add_heading(doc, "2.2. User Stories", level=2)
    add_para(doc,
        "Cada historia sigue el formato Scrum: Como [rol] quiero [funcionalidad] para [beneficio], "
        "acompañada de criterios de aceptación verificables.")

    user_stories = [
        ("US-01", "Registro de socio nuevo desde formulario web",
         "Como visitante quiero registrarme como nuevo socio ingresando mis datos personales para acceder a los servicios de GymMax.",
         ["El formulario solicita: nombres, apellidos, DNI (8 dígitos), correo, celular (9 dígitos), contraseña (mín 8 caracteres) y aceptación de términos.",
          "El DNI y el correo deben ser únicos en el sistema.",
          "La contraseña debe tener al menos una mayúscula, minúscula y un número.",
          "Al enviar datos válidos, el socio queda registrado y es redirigido al dashboard.",
          "Se almacena en las tablas USUARIO y SOCIO."]),

        ("US-02", "Inicio de sesión con correo y contraseña",
         "Como usuario registrado quiero iniciar sesión para acceder a las funcionalidades de mi rol.",
         ["Formulario con correo, contraseña, checkbox 'Recordarme' y enlace 'Olvidaste tu contraseña'.",
          "Si credenciales válidas y rol SOCIO, redirige al dashboard del socio.",
          "Si credenciales válidas y rol ADMIN, redirige al dashboard administrativo.",
          "Si inválidas, muestra 'Correo o contraseña incorrectos'.",
          "Después de 3 intentos fallidos, muestra advertencia.",
          "Con 'Recordarme', la sesión persiste por 7 días."]),

        ("US-03", "Recuperar contraseña olvidada",
         "Como usuario que olvidó su contraseña quiero recuperarla mediante correo electrónico para volver a acceder.",
         ["El usuario ingresa su correo en la pantalla de recuperación.",
          "Si el correo existe, se envía un enlace temporal (válido por 30 minutos).",
          "Al hacer clic, el usuario puede establecer una nueva contraseña.",
          "La nueva contraseña cumple las reglas de seguridad del registro.",
          "Si el correo no existe, mensaje genérico (por seguridad)."]),

        ("US-04", "Visualización del dashboard del socio",
         "Como socio quiero ver mi dashboard personal para conocer el estado de mi membresía y acciones disponibles.",
         ["Saludo personalizado 'Hola, [Nombre]'.",
          "Tarjeta con membresía activa: nombre del plan y vencimiento.",
          "Contadores de asistencias y reservas activas.",
          "Próxima clase reservada con nombre, hora y sede.",
          "Accesos rápidos a Ver sedes, Reservar clase, Mis reservas.",
          "Alerta visual si membresía vencida o por vencer (≤7 días)."]),

        ("US-05", "Consulta del listado de sedes",
         "Como socio quiero ver todas las sedes para elegir cuál visitar.",
         ["Lista o grid con todas las sedes activas.",
          "Cada sede muestra: nombre, dirección, distrito, horario y cupos disponibles (con color).",
          "Al hacer clic, navega al detalle con clases y horarios.",
          "Campo de búsqueda para filtrar por distrito o nombre."]),

        ("US-06", "Filtrar sedes por distrito",
         "Como socio quiero filtrar sedes por distrito para encontrar la más cercana.",
         ["Campo de búsqueda o dropdown con todos los distritos.",
          "Filtro dinámico sin recargar página.",
          "Si no hay sedes, mensaje 'No se encontraron sedes en este distrito'.",
          "Botón 'Limpiar filtros' para resetear."]),

        ("US-07", "Ver clases disponibles en una sede",
         "Como socio quiero ver las clases de una sede para decidir cuál reservar.",
         ["Nombre de la sede en el encabezado.",
          "Selector de fecha (por defecto hoy).",
          "Listado de clases del día: nombre, hora, instructor, tipo y cupos (ej. 10/15).",
          "Botón 'Reservar' si hay cupos o 'Sin cupos' si está llena.",
          "Estado de cupos se actualiza al refrescar."]),

        ("US-08", "Reservar una clase",
         "Como socio con membresía activa quiero reservar una clase para asegurar mi cupo.",
         ["El sistema verifica membresía activa antes de reservar.",
          "Verifica que haya cupos disponibles.",
          "Si todo válido, reserva con estado CONFIRMADA y descuenta un cupo.",
          "Si no hay cupos, ofrece 'Lista de espera'.",
          "No permite reservar la misma clase dos veces.",
          "Mensaje de confirmación con detalles de la reserva."]),

        ("US-09", "Ver mis reservas",
         "Como socio quiero ver mis reservas (próximas y pasadas) para controlar mis clases.",
         ["Dos pestañas: 'Próximas (N)' y 'Pasadas'.",
          "Próximas: clase, fecha, hora, sede, estado.",
          "Pasadas: clase, fecha, sede y si asistió o no.",
          "Botón cancelar en las próximas.",
          "Botón '+ Nueva reserva'."]),

        ("US-10", "Cancelar una reserva",
         "Como socio quiero cancelar una reserva confirmada para liberar el cupo.",
         ["Solo cancelables las CONFIRMADAS futuras.",
          "Mínimo 2 horas de anticipación.",
          "Si cumple, estado pasa a CANCELADA y cupo se devuelve.",
          "Si no, mensaje 'No se puede cancelar (menos de 2h)'.",
          "Modal de confirmación antes de cancelar."]),

        ("US-11", "Contratar o renovar plan de membresía",
         "Como socio quiero contratar o renovar mi plan para mantener acceso al gimnasio.",
         ["Planes disponibles con precios, beneficios y duración.",
          "Checkout con resumen del pedido.",
          "Si tiene membresía activa, se renueva (suma duración al vencimiento).",
          "Si no tiene, fecha de inicio es hoy.",
          "Selecciona método de pago (Yape/Plin/Tarjeta).",
          "Al confirmar, registra en MEMBRESIA y PAGO.",
          "Genera comprobante con número de operación."]),

        ("US-12", "Visualizar historial de pagos",
         "Como socio quiero ver mis pagos para llevar control de mis gastos.",
         ["Lista con fecha, plan, método, número operación, monto y estado.",
          "Filtros por rango de fechas y estado.",
          "Enlace para descargar comprobante PDF.",
          "Ordenados del más reciente al más antiguo."]),

        ("US-13", "Login del administrador",
         "Como administrador quiero iniciar sesión y ser redirigido a mi dashboard.",
         ["Mismo formulario de login que socio.",
          "Sistema valida rol ADMIN tras autenticar.",
          "Si ADMIN, redirige a adminDashboard.jsp.",
          "Si SOCIO, redirige a dashboardSocio.jsp.",
          "Socio no puede acceder a rutas administrativas."]),

        ("US-14", "Dashboard administrativo con KPIs",
         "Como administrador quiero ver indicadores clave para tomar decisiones.",
         ["4 tarjetas KPI: socios activos, ingresos del mes, vencimientos próximos, reservas del día.",
          "Cada KPI muestra variación % vs período anterior.",
          "Gráfico de líneas con ingresos últimos 6 meses.",
          "Lista de actividad reciente (últimas 5-10 acciones).",
          "Datos reales de la base de datos."]),

        ("US-15", "CRUD completo de socios",
         "Como administrador quiero gestionar socios para mantener el registro actualizado.",
         ["Tabla paginada con ID, Nombre, DNI, Correo, Plan, Sede, Estado, Vencimiento, Acciones.",
          "Filtros por nombre/DNI/correo, sede, plan, estado.",
          "Botón 'Nuevo socio' abre formulario.",
          "Acciones: editar y eliminar con confirmación.",
          "Eliminar es soft delete (marca como inactivo).",
          "Cambios persisten en MySQL."]),

        ("US-16", "CRUD completo de sedes",
         "Como administrador quiero gestionar sedes para administrar la cadena.",
         ["Listado con nombre, dirección, distrito, horarios, capacidad, estado.",
          "Botón 'Nueva sede' abre formulario.",
          "Editar y eliminar (no se puede eliminar con socios/reservas activas).",
          "Al eliminar con dependencias, se ofrece desactivar.",
          "Cambios reflejados en el listado al socio."]),

        ("US-17", "Generar reportes",
         "Como administrador quiero generar reportes filtrables para analizar el negocio.",
         ["Selecciona tipo: Ingresos por sede, Membresías por estado, Reservas por clase, Morosidad.",
          "Filtros por rango de fechas y sede.",
          "Gráfico (barras o líneas) según tipo.",
          "Resumen numérico con tabla.",
          "Botones Generar, Resetear, Exportar Excel, Exportar PDF."]),

        ("US-18", "Exportar reportes a Excel y PDF",
         "Como administrador quiero exportar reportes para compartirlos o archivarlos.",
         ["Botón 'Excel' descarga .xlsx con datos actuales.",
          "Botón 'PDF' descarga .pdf con formato profesional.",
          "Nombre del archivo incluye tipo y fecha.",
          "Respeta los filtros aplicados al momento de exportar."]),
    ]

    for code, titulo, desc, criterios in user_stories:
        add_para(doc, f"{code} — {titulo}", bold=True, size=12, color=ORANGE)
        add_para(doc, "Descripción: ", bold=True)
        add_para(doc, desc)
        add_para(doc, "Criterios de aceptación:", bold=True)
        for c in criterios:
            add_bullet(doc, c)
        doc.add_paragraph()

    add_page_break(doc)

    # 2.3 Flujos
    add_heading(doc, "2.3. Diagrama de flujo de navegación entre Wireframes", level=2)
    add_para(doc, "Flujo del Socio:", bold=True)
    add_para(doc,
        "El socio ingresa al login → si sus credenciales son válidas accede al dashboard; "
        "si no, puede registrarse. Desde el dashboard puede navegar a sedes, reservar clases "
        "y ver sus reservas, terminando con logout.")
    add_para(doc, "[INSERTAR AQUÍ captura del diagrama de flujo del Socio del Avance 01]",
             italic=True, color=GRAY)
    add_para(doc, "Flujo del Administrador:", bold=True)
    add_para(doc,
        "El admin se autentica → el sistema valida rol → accede al dashboard administrativo "
        "desde donde gestiona socios, sedes y genera reportes, terminando con logout.")
    add_para(doc, "[INSERTAR AQUÍ captura del diagrama de flujo del Admin del Avance 01]",
             italic=True, color=GRAY)
    add_page_break(doc)

    # 2.4 Wireframes
    add_heading(doc, "2.4. Wireframes — Responsive Web Design", level=2)
    add_heading(doc, "2.4.1. Versión Mobile", level=3)
    add_para(doc, "[INSERTAR AQUÍ capturas de los 6 mockups móviles del Avance 01: "
                  "login, registro, dashboard socio, listado sedes, reservar clase, mis reservas]",
             italic=True, color=GRAY)
    add_heading(doc, "2.4.2. Versión Desktop", level=3)
    add_para(doc, "[INSERTAR AQUÍ capturas de los mockups desktop del Avance 01: "
                  "dashboard admin, gestión de socios, generación de reportes]",
             italic=True, color=GRAY)
    add_page_break(doc)

    # 2.5 BD
    add_heading(doc, "2.5. Modelo de Base de Datos", level=2)
    add_heading(doc, "2.5.1. Modelo Lógico", level=3)
    add_para(doc,
        "El modelo lógico representa las entidades del negocio sin atarse a un SGBD específico. "
        "Está compuesto por 10 entidades: Usuario, Socio, Plan, Membresía, Pago, Sede, Clase, "
        "Horario_Clase, Reserva, Asistencia, con 10 relaciones entre ellas.")
    add_para(doc, "[INSERTAR AQUÍ captura del diagrama lógico hecho en Lucid Chart]",
             italic=True, color=GRAY)

    add_heading(doc, "2.5.2. Modelo Físico (MySQL 8.0)", level=3)
    add_para(doc,
        "El modelo físico implementa el lógico en MySQL 8.0 con tipos de dato específicos "
        "(VARCHAR, INT, DATE, ENUM), claves primarias y foráneas, restricciones (UNIQUE, "
        "NOT NULL, CHECK), índices para optimización y motor InnoDB para soporte transaccional.")
    add_para(doc, "[INSERTAR AQUÍ captura del diagrama físico hecho en Lucid Chart]",
             italic=True, color=GRAY)
    add_para(doc, "Las tablas implementadas son:")
    tablas_db = [
        ("USUARIO", "Credenciales y datos generales del usuario (rol Socio o Admin)"),
        ("SOCIO", "Datos específicos del socio (DNI, celular, fecha nacimiento, etc.)"),
        ("PLAN_MEMBRESIA", "Catálogo de planes (Básico, Premium, Anual)"),
        ("MEMBRESIA", "Contrato vigente entre un socio y un plan"),
        ("PAGO", "Transacciones financieras asociadas a una membresía"),
        ("SEDE", "Locales físicos del gimnasio"),
        ("CLASE", "Actividades programadas en una sede"),
        ("HORARIO_CLASE", "Instancias específicas de una clase en una fecha"),
        ("RESERVA", "Cupo tomado por un socio en una clase y horario"),
        ("ASISTENCIA", "Registro de ingreso/salida del socio a una sede"),
    ]
    add_table_from_data(doc, ["Tabla", "Descripción"], tablas_db)

    add_page_break(doc)

    # 2.6 Tecnologías
    add_heading(doc, "2.6. Resumen de tecnologías aplicadas", level=2)

    add_heading(doc, "2.6.1. Lenguajes y frameworks", level=3)
    tech_lang = [
        ("Java", "11 (LTS)", "Backend", "Lenguaje principal del sistema"),
        ("Jakarta EE", "10.0.0", "Backend", "Plataforma empresarial (Servlets, JSP, JDBC, CDI)"),
        ("JSP", "3.1", "Frontend", "Páginas dinámicas server-side para las vistas"),
        ("HTML5", "—", "Frontend", "Estructura semántica de las páginas"),
        ("CSS3", "—", "Frontend", "Estilos y responsive"),
        ("JavaScript", "ES6+", "Frontend", "Interactividad y validación cliente"),
        ("Bootstrap", "5.3", "Frontend", "Framework CSS Mobile First"),
        ("Bootstrap Icons", "1.11.3", "Frontend", "Librería de íconos"),
        ("JSTL", "3.0", "Frontend", "Etiquetas estándar para JSP"),
    ]
    add_table_from_data(doc, ["Tecnología", "Versión", "Capa", "Propósito"], tech_lang)

    add_heading(doc, "2.6.2. Base de datos y persistencia", level=3)
    tech_db = [
        ("MySQL Community", "8.0", "BD", "Sistema gestor relacional"),
        ("MySQL Workbench", "8.0", "Herramienta", "Cliente gráfico para BD"),
        ("JDBC", "4.3", "Backend", "API estándar de Java para BD"),
        ("MySQL Connector/J", "8.0.33", "Backend", "Driver JDBC oficial de MySQL"),
    ]
    add_table_from_data(doc, ["Tecnología", "Versión", "Capa", "Propósito"], tech_db)

    add_heading(doc, "2.6.3. Servidor y despliegue", level=3)
    tech_srv = [
        ("Apache Tomcat", "10.1.x", "Servidor", "Contenedor Servlet/JSP compatible Jakarta EE 10"),
        ("Apache Maven", "3.9.x", "Build", "Gestión de dependencias y empaquetado WAR"),
    ]
    add_table_from_data(doc, ["Tecnología", "Versión", "Capa", "Propósito"], tech_srv)

    add_heading(doc, "2.6.4. Patrones arquitectónicos y de diseño", level=3)
    patrones = [
        ("Arquitectura en capas", "Arquitectura", "Separación presentación/controlador/fachada/DAO/datos"),
        ("MVC", "Arquitectura", "JSP=Vista, Servlet=Controller, DTO=Model"),
        ("DAO", "Diseño", "Aislamiento de acceso a datos con interfaces + impl"),
        ("Facade", "Diseño", "Encapsulamiento de lógica de negocio"),
        ("DTO", "Diseño", "Objetos planos para transportar datos entre capas"),
        ("Singleton", "Diseño", "Helper de conexión a BD único"),
    ]
    add_table_from_data(doc, ["Patrón", "Categoría", "Propósito en el sistema"], patrones)

    add_heading(doc, "2.6.5. Metodología y herramientas de gestión", level=3)
    meto = [
        ("Scrum", "Metodología", "Marco ágil con sprints de 2 semanas"),
        ("Trello/Jira", "Gestión", "Tablero Kanban del Product Backlog"),
        ("Lucid Chart", "Modelado", "Diagramas de flujo, BD lógico/físico, clases UML"),
        ("Microsoft Word", "Documentación", "Informes (APA 7ma edición)"),
        ("Microsoft PowerPoint", "Presentación", "Diapositivas de avances"),
        ("Git + GitHub", "Versionamiento", "Repositorio del código fuente del grupo"),
    ]
    add_table_from_data(doc, ["Herramienta", "Categoría", "Propósito"], meto)

    add_heading(doc, "2.6.6. Entorno de desarrollo", level=3)
    dev = [
        ("Apache NetBeans", "21+", "IDE", "Entorno integrado para Java EE"),
        ("OpenJDK", "11+", "Runtime", "Kit de desarrollo Java"),
        ("Postman", "—", "API testing", "Testing de endpoints"),
        ("Chrome/Edge", "Última", "Navegador", "Visualización y pruebas frontend"),
        ("DevTools", "—", "Debugging", "Inspección de elementos y red"),
    ]
    add_table_from_data(doc, ["Tecnología", "Versión", "Categoría", "Propósito"], dev)

    add_page_break(doc)

    # 2.7 Enlaces
    add_heading(doc, "2.7. Enlaces a artefactos", level=2)
    add_para(doc, "Los artefactos del proyecto están disponibles en los siguientes recursos:")
    add_bullet(doc, "[INSERTAR ENLACE Lucid Chart con todos los diagramas]", "Diagramas en Lucid Chart:")
    add_bullet(doc, "[INSERTAR ENLACE GitHub del repositorio del grupo]", "Repositorio Git:")
    add_bullet(doc, "[INSERTAR ENLACE de carpeta compartida con Word, PPT, mockups]", "Drive del proyecto:")
    add_bullet(doc, "[INSERTAR ENLACE Trello/Jira con el Sprint Backlog]", "Tablero del backlog:")

    add_page_break(doc)


# ----- CAPÍTULO III -----
def build_cap3(doc):
    add_heading(doc, "Capítulo III: Backend", level=1, color=ORANGE)

    add_heading(doc, "3.1. Diagrama de Clases del Sistema (UML)", level=2)
    add_para(doc,
        "El diagrama de clases del sistema GymMax se organiza bajo la arquitectura por capas "
        "MVC + DAO + Facade. La estructura completa cuenta con 39 clases distribuidas en "
        "5 paquetes:")
    add_bullet(doc, "Singleton helper de conexión a MySQL.", "com.gymmax (raíz):")
    add_bullet(doc, "5 Servlets que actúan como controladores (Auth, Socio, Reserva, Membresia, Admin).",
               "com.gymmax.controller:")
    add_bullet(doc, "9 interfaces DAO (I*DAO) + 9 implementaciones DAO (*DAOImpl).",
               "com.gymmax.dao:")
    add_bullet(doc, "10 clases POJO (Usuario, Socio, Plan, Membresia, Pago, Sede, Clase, HorarioClase, Reserva, Asistencia).",
               "com.gymmax.dto:")
    add_bullet(doc, "5 Facades que orquestan reglas de negocio (Auth, Socio, Reserva, Membresia, Admin).",
               "com.gymmax.facade:")
    add_para(doc, "[INSERTAR AQUÍ captura del diagrama de clases UML hecho en Lucid Chart]",
             italic=True, color=GRAY)
    add_para(doc,
        "Sugerencia: dividir el diagrama en 3 partes para mayor legibilidad:")
    add_bullet(doc, "Las 10 clases DTO con sus asociaciones.", "Diagrama 1 (DTOs):")
    add_bullet(doc, "Las 9 interfaces y sus 9 implementaciones DAO.", "Diagrama 2 (DAOs):")
    add_bullet(doc, "Las 5 Facades y los 5 Servlets, con sus dependencias.", "Diagrama 3 (Facades + Controllers):")

    add_page_break(doc)

    # 3.2 Diccionario
    add_heading(doc, "3.2. Diccionario de Clases", level=2)
    add_para(doc,
        "A continuación se describe cada clase del sistema con sus atributos principales, "
        "tipos de dato y responsabilidades.")

    # Diccionario condensado (atributos clave)
    dto_data = [
        ("Usuario",
         "idUsuario:int, correo:String, password:String, nombres:String, apellidos:String, rol:String, creadoEn:LocalDateTime",
         "Persona registrada con credenciales de acceso (rol SOCIO o ADMIN)."),
        ("Socio",
         "idSocio:int, idUsuario:int, dni:String, celular:String, fechaNac:LocalDate, genero:char, direccion:String, fechaReg:LocalDate",
         "Cliente del gimnasio (extiende a Usuario con datos personales)."),
        ("Plan",
         "idPlan:int, nombre:String, tipo:String, precio:double, duracionDias:int, activo:boolean",
         "Tipo de membresía ofrecido (Básico, Premium, Anual)."),
        ("Membresia",
         "idMembresia:int, idSocio:int, idPlan:int, fechaInicio:LocalDate, fechaFin:LocalDate, estado:String, monto:double, renovacionAuto:boolean",
         "Contrato vigente entre socio y plan. Métodos: isVigente(), diasParaVencer()."),
        ("Pago",
         "idPago:int, idMembresia:int, metodo:String, monto:double, fechaPago:LocalDateTime, nroOperacion:String, estado:String, comprobanteUrl:String",
         "Transacción financiera asociada a una membresía."),
        ("Sede",
         "idSede:int, nombre:String, direccion:String, distrito:String, telefono:String, horaApertura:LocalTime, horaCierre:LocalTime, capacidad:int",
         "Local físico del gimnasio."),
        ("Clase",
         "idClase:int, idSede:int, nombre:String, tipo:String, instructor:String, diaSemana:int, horaInicio:LocalTime, cupoMaximo:int",
         "Actividad programada en una sede (CrossFit, Yoga, etc.)."),
        ("HorarioClase",
         "idHorario:int, idClase:int, fechaEspecifica:LocalDate, cupoActual:int, cupoDisponible:int, estado:String",
         "Instancia específica de una clase en fecha concreta. Método: hayCupo()."),
        ("Reserva",
         "idReserva:int, idSocio:int, idClase:int, fecha:LocalDate, hora:LocalTime, estado:String, creadoEn:LocalDateTime, canceladoEn:LocalDateTime",
         "Cupo tomado por un socio. Método: puedeCancelarse()."),
        ("Asistencia",
         "idAsistencia:int, idSocio:int, idSede:int, fechaHora:LocalDateTime, tipo:String, registradoPor:int",
         "Registro de ingreso o salida del socio a una sede."),
    ]
    add_heading(doc, "3.2.1. Capa DTO (com.gymmax.dto)", level=3)
    add_table_from_data(doc, ["Clase", "Atributos", "Descripción / Métodos"], dto_data)

    add_heading(doc, "3.2.2. Capa DAO (com.gymmax.dao)", level=3)
    add_para(doc,
        "Cada entidad principal cuenta con una interfaz I*DAO que define el contrato CRUD "
        "y una implementación *DAOImpl que usa JDBC + Conexion.getConnection().")
    dao_data = [
        ("IUsuarioDAO / UsuarioDAOImpl", "CRUD de Usuario + búsqueda por correo."),
        ("ISocioDAO / SocioDAOImpl", "CRUD de Socio + búsqueda por DNI/usuario + filtros."),
        ("IPlanDAO / PlanDAOImpl", "CRUD de Plan + listar activos + desactivar."),
        ("IMembresiaDAO / MembresiaDAOImpl", "CRUD + buscar activa por socio + listar por vencer."),
        ("IPagoDAO / PagoDAOImpl", "CRUD + listar por socio/membresía + sumar ingresos."),
        ("ISedeDAO / SedeDAOImpl", "CRUD + filtrar por distrito."),
        ("IClaseDAO / ClaseDAOImpl", "CRUD + listar por sede + por fecha."),
        ("IReservaDAO / ReservaDAOImpl", "CRUD + listar por socio/clase + cancelar + validar duplicado."),
        ("IAsistenciaDAO / AsistenciaDAOImpl", "Registrar + listar por socio + contar por sede."),
    ]
    add_table_from_data(doc, ["Interfaz / Implementación", "Responsabilidad principal"], dao_data)

    add_heading(doc, "3.2.3. Capa Facade (com.gymmax.facade)", level=3)
    facade_data = [
        ("AuthFacade",
         "login(correo, password), registrarSocio(socio, password), recuperarPassword(correo), cambiarPassword(idUsuario, nueva)"),
        ("SocioFacade",
         "obtenerDashboard(idSocio), listarSedes(), listarSedesPorDistrito(distrito), listarClasesPorSedeFecha(idSede, fecha)"),
        ("ReservaFacade",
         "reservar(idSocio, idClase, fecha, hora) — valida membresía+cupos+duplicado | cancelar(idReserva) — regla 2h | listarMisReservas(idSocio, soloFuturas)"),
        ("MembresiaFacade",
         "contratar(idSocio, idPlan, metodoPago) — crea Membresia+Pago en transacción | renovar(idMembresia) | listarHistorialPagos(idSocio)"),
        ("AdminFacade",
         "obtenerDashboardAdmin() — KPIs | generarReporteIngresos(desde, hasta, idSede) | exportarExcel(reporte) | exportarPDF(reporte)"),
    ]
    add_table_from_data(doc, ["Facade", "Métodos públicos (con reglas de negocio)"], facade_data)

    add_heading(doc, "3.2.4. Capa Controller (com.gymmax.controller)", level=3)
    ctrl_data = [
        ("AuthServlet (@WebServlet(\"/AuthServlet\"))",
         "doGet: muestra login/registro. doPost: procesa login o registro, setea sesión por rol."),
        ("SocioServlet (@WebServlet(\"/SocioServlet\"))",
         "doGet: dashboard, listado de sedes o detalle. Usa AuthFacade para verificar sesión."),
        ("ReservaServlet (@WebServlet(\"/ReservaServlet\"))",
         "doGet: muestra formulario o 'Mis reservas'. doPost: procesa reserva o cancelación."),
        ("MembresiaServlet (@WebServlet(\"/MembresiaServlet\"))",
         "doGet: muestra planes. doPost: procesa contratación/renovación."),
        ("AdminServlet (@WebServlet(\"/AdminServlet\"))",
         "doGet: dashboard, gestión, reportes (según parámetro view). doPost: CRUDs y reportes."),
    ]
    add_table_from_data(doc, ["Servlet (mapping)", "Responsabilidad"], ctrl_data)

    add_heading(doc, "3.2.5. Clase utilitaria", level=3)
    add_para(doc, "Conexion (com.gymmax.Conexion)", bold=True)
    add_para(doc,
        "Singleton que centraliza la creación de conexiones JDBC a MySQL. "
        "Atributos estáticos: URL, USER, PASSWORD. "
        "Método principal: getConnection() que retorna java.sql.Connection.")

    add_heading(doc, "3.2.6. Resumen del diccionario", level=3)
    resumen = [
        ("DTO", "10 clases"),
        ("DAO interfaces", "9"),
        ("DAO implementaciones", "9"),
        ("Facades", "5"),
        ("Controllers (Servlets)", "5"),
        ("Utilitaria (Conexion)", "1"),
        ("TOTAL", "39"),
    ]
    add_table_from_data(doc, ["Capa", "Cantidad de clases/interfaces"], resumen)


def build_word():
    print("Generando Word...")
    doc = Document()

    # Márgenes
    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin = Cm(3)
        section.right_margin = Cm(3)

    build_caratula(doc)
    build_indice(doc)
    build_cap1(doc)
    build_cap2(doc)
    build_cap3(doc)

    doc.save(OUT_DOCX)
    print(f"  -> {OUT_DOCX}")


# ============================================================
#               GENERACIÓN DEL POWERPOINT (.pptx)
# ============================================================

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(prs, title, bullets, notas=""):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            tf.text = b
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            p.text = b
        p.font.size = PPt(18)
    if notas:
        slide.notes_slide.notes_text_frame.text = notas
    return slide


def build_pptx():
    print("Generando PowerPoint...")
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)

    # Slide 1: Portada
    add_title_slide(prs,
        "GymMax - Sistema Integral de Gestión de Gimnasios",
        "Avance 02 - Diseño y Arquitectura\nUTP - Desarrollo Web Integrado - Sec. 27672\nDocente: Juan Manuel Rodríguez del Águila\nMayo 2026")

    # Slide 2: Integrantes
    add_content_slide(prs, "Integrantes del Grupo 3", [
        "Choque Anchante, Niurka Yasbeth",
        "Diaz Culqui, Nehemias",
        "Purizaca Ipanaque, Dennys Marlon",
        "Soria Chavez, Ian",
        "Torre Escobar, Oliver",
        "Valladolid Llenque, Alexander",
    ])

    # Slide 3: Agenda
    add_content_slide(prs, "Agenda", [
        "1. Introducción y problemática",
        "2. Objetivos y alcance",
        "3. Características del producto",
        "4. Product Backlog y User Stories",
        "5. Wireframes y flujos de navegación",
        "6. Modelo de Base de Datos (Lógico + Físico)",
        "7. Diagrama de Clases UML",
        "8. Tecnologías utilizadas",
        "9. Próximos pasos",
    ], "En los próximos 15 minutos recorreremos los 9 puntos que definen el diseño completo del sistema GymMax.")

    # Slide 4: El producto
    add_content_slide(prs, "¿Qué es GymMax?", [
        "Plataforma web centralizada para gestión de cadenas de gimnasios",
        "Construida sobre Java EE + MySQL + Bootstrap 5.3",
        "Enfoque Mobile First: accesible desde cualquier dispositivo",
        "Centraliza: sedes, socios, membresías, reservas, pagos y reportes",
    ])

    # Slide 5: Problemática
    add_content_slide(prs, "¿Qué problemas resolvemos?", [
        "Pérdida de ingresos por falta de control de vencimientos",
        "Sobrecarga administrativa y procesos manuales",
        "Colas largas en recepción en horarios pico",
        "Mala experiencia del socio: no puede autogestionarse desde el celular",
        "Sistemas existentes lentos y poco escalables",
    ], "Estos 5 problemas fueron validados en investigación de antecedentes del sector peruano (SmartGym, GymSys, BodyTrack).")

    # Slide 6: Objetivos
    add_content_slide(prs, "Objetivo general y específicos", [
        "GENERAL: Desarrollar aplicación web JEE que centralice sedes, socios, membresías y reservas con interfaz Mobile First.",
        "1. Módulo de consulta de sedes con cupos por clase",
        "2. Formulario responsive de registro con validación tiempo real",
        "3. Sistema de planes (Básico/Premium/Anual) con renovación autónoma",
        "4. Diseño Mobile First con Bootstrap 5.3",
        "5. Panel admin con KPIs y reportes exportables",
        "6. Arquitectura por capas MVC + DAO + Facade",
    ])

    # Slide 7: Alcance
    add_content_slide(prs, "Alcance del proyecto", [
        "SÍ INCLUYE: Multi-sede, multi-rol, membresías digitales, dashboard KPIs, reportes",
        "NO INCLUYE: pagos en efectivo, app móvil nativa, biometría, marketing masivo",
        "PÚBLICO OBJETIVO: cadenas peruanas de 2-10 sedes",
        "ARQUITECTURA: Jakarta EE con alta disponibilidad y escalabilidad",
    ])

    # Slide 8: Características
    add_content_slide(prs, "10 características clave del producto", [
        "1. Multi-sede centralizada",
        "2. Sistema de roles Socio/Admin",
        "3. Mobile First responsive",
        "4. Pagos digitales (Yape/Plin/Tarjeta)",
        "5. Dashboard administrativo con KPIs",
        "6. Reservas online con control de cupos",
        "7. Validación en tiempo real",
        "8. Reportes exportables a Excel/PDF",
        "9. Seguridad y control de sesión por rol",
        "10. Arquitectura por capas MVC + DAO + Facade",
    ])

    # Slide 9: Product Backlog
    add_content_slide(prs, "Product Backlog - Requerimientos", [
        "28 Requerimientos Funcionales + 12 No Funcionales",
        "Distribuidos en 5 sprints (134 Story Points totales)",
        "Sprint 1: Auth + Socio base - 22 SP",
        "Sprint 2: Sedes + Reservas - 37 SP",
        "Sprint 3: Membresías + Pagos - 41 SP",
        "Sprint 4: Reportes + Asistencias - 26 SP",
        "Sprint 5: Notificaciones - 8 SP",
    ], "El backlog completo está en el informe Word con las 40 historias detalladas.")

    # Slide 10: User Stories destacadas
    add_content_slide(prs, "User Stories destacadas", [
        "US-08 Reservar clase: como socio con membresía activa quiero reservar una clase para asegurar mi cupo (verifica membresía, descuenta cupo, confirma)",
        "US-14 Dashboard Admin: como admin quiero ver KPIs en tiempo real para decidir (4 KPIs, variación %, gráfico 6 meses)",
        "US-17 Generar reportes: como admin quiero reportes filtrables (4 tipos, filtros, exportar Excel/PDF)",
    ], "El documento incluye las 18 historias completas con criterios de aceptación.")

    # Slide 11: Flujo Socio
    add_content_slide(prs, "Diagrama de flujo - Socio", [
        "INSERTAR aquí captura del diagrama del Avance 01",
        "Login -> validar credenciales -> Dashboard del socio",
        "Desde dashboard: Ver sedes / Reservar clase / Mis reservas",
        "Si no tiene cuenta: ir a Registro",
    ], "El socio ingresa al login; si sus credenciales son válidas accede al dashboard; si no, puede registrarse.")

    # Slide 12: Flujo Admin
    add_content_slide(prs, "Diagrama de flujo - Administrador", [
        "INSERTAR aquí captura del diagrama del Avance 01",
        "Login -> validar rol ADMIN -> Dashboard administrativo",
        "Desde dashboard: Gestión socios / Gestión sedes / Reportes",
    ], "El admin se autentica, el sistema valida el rol y accede al panel administrativo.")

    # Slide 13: Mockups Mobile (1)
    add_content_slide(prs, "Wireframes Mobile (1/2)", [
        "INSERTAR aquí los mockups del Avance 01:",
        "Login del socio",
        "Registro de socio",
        "Dashboard del socio",
    ])

    # Slide 14: Mockups Mobile (2)
    add_content_slide(prs, "Wireframes Mobile (2/2)", [
        "INSERTAR aquí los mockups del Avance 01:",
        "Listado de sedes",
        "Reservar clase",
        "Mis reservas",
    ])

    # Slide 15: Mockups Desktop
    add_content_slide(prs, "Wireframes Desktop - Admin", [
        "INSERTAR aquí mockups del Avance 01:",
        "Dashboard administrativo con KPIs",
        "Gestión de socios (tabla CRUD)",
    ])

    # Slide 16: Reportes Desktop
    add_content_slide(prs, "Wireframes Desktop - Reportes", [
        "INSERTAR aquí mockup del Avance 01:",
        "Generación de reportes",
        "Gráfico de barras + tabla de resumen numérico",
    ])

    # Slide 17: BD Lógico
    add_content_slide(prs, "Modelo de Base de Datos - LÓGICO", [
        "INSERTAR aquí captura del diagrama lógico de Lucid Chart",
        "10 entidades, 10 relaciones, normalizado a 3FN",
        "Entidades: Usuario, Socio, Plan, Membresía, Pago, Sede, Clase, HorarioClase, Reserva, Asistencia",
    ], "El modelo lógico representa las entidades del negocio sin atarse a un SGBD específico.")

    # Slide 18: BD Físico
    add_content_slide(prs, "Modelo de Base de Datos - FÍSICO", [
        "INSERTAR aquí captura del diagrama físico de Lucid Chart",
        "Motor InnoDB con integridad referencial",
        "Charset utf8mb4 (soporte de caracteres especiales)",
        "Script SQL listo (en anexos del Word)",
    ])

    # Slide 19: Tecnologías
    add_content_slide(prs, "Stack tecnológico de GymMax", [
        "Lenguaje: Java 11 + Jakarta EE 10",
        "Frontend: JSP + Bootstrap 5.3 + HTML5",
        "Controlador: Servlets (Jakarta 6.0)",
        "Persistencia: JDBC + DAO + MySQL Connector",
        "Base de datos: MySQL 8.0 Community",
        "Servidor: Apache Tomcat 10.1",
        "IDE: NetBeans 21",
        "Diagramas: Lucid Chart",
        "Metodología: Scrum (sprints de 2 semanas)",
    ])

    # Slide 20: UML 1
    add_content_slide(prs, "Diagrama de Clases UML (1/2)", [
        "INSERTAR aquí captura UML (DTOs + DAOs)",
        "10 DTOs (entidades del modelo de datos)",
        "9 interfaces DAO + 9 implementaciones DAO",
        "Programación a interfaces para desacoplar persistencia",
    ], "Cada entidad tiene su DTO, una interfaz DAO que define el contrato CRUD, y una implementación que usa JDBC.")

    # Slide 21: UML 2
    add_content_slide(prs, "Diagrama de Clases UML (2/2)", [
        "INSERTAR aquí captura UML (Facades + Controllers)",
        "5 Facades (Auth, Socio, Reserva, Membresia, Admin)",
        "5 Servlets (uno por Facade)",
        "1 Conexion (singleton helper de JDBC)",
        "Facade encapsula reglas de negocio; Servlets coordinan MVC",
    ])

    # Slide 22: Diccionario
    add_content_slide(prs, "Diccionario de Clases - Resumen", [
        "DTO: 10 clases",
        "DAO interfaces: 9",
        "DAO implementaciones: 9",
        "Facades: 5",
        "Controllers (Servlets): 5",
        "Utilitaria (Conexion): 1",
        "TOTAL: 39 clases",
    ], "El diccionario completo con atributos, métodos y descripciones está en el informe Word.")

    # Slide 23: Patrones
    add_content_slide(prs, "Patrones de diseño aplicados", [
        "MVC: Modelo (DTO) - Vista (JSP) - Controlador (Servlet)",
        "DAO: aislamiento del acceso a datos",
        "Facade: encapsulamiento de reglas de negocio",
        "Singleton: helper Conexion único para JDBC",
        "DTO: objetos planos para transportar datos entre capas",
    ])

    # Slide 24: Próximos pasos
    add_content_slide(prs, "Hoja de ruta hacia la entrega final", [
        "Sprint 1-2 (Sem 11-12): Autenticación + dashboard socio + listado de sedes",
        "Sprint 3 (Sem 13): Reservas de clases + membresías y pagos",
        "Sprint 4 (Sem 14): Panel administrativo + reportes + exportación",
        "Sprint 5 (Sem 15): Notificaciones + pruebas integrales + documentación final",
    ])

    # Slide 25: Conclusiones
    add_content_slide(prs, "Conclusiones del Avance 02", [
        "Diseño completo del sistema cubierto: 18 user stories, 40 requerimientos",
        "Modelo de datos validado: 10 tablas normalizadas + script SQL listo",
        "Arquitectura por capas definida: 39 clases organizadas en 5 paquetes",
        "Stack tecnológico consolidado: Java EE + MySQL + Bootstrap + Tomcat",
    ])

    # Slide 26: Gracias
    add_title_slide(prs, "¡Gracias!", "¿Preguntas?\n\nEquipo Grupo 3 - UTP 2026")

    prs.save(OUT_PPTX)
    print(f"  -> {OUT_PPTX}")


# ============================================================
if __name__ == "__main__":
    build_word()
    build_pptx()
    print("\n=== Generación completa ===")
    print(f"Word:       {OUT_DOCX}")
    print(f"PowerPoint: {OUT_PPTX}")
