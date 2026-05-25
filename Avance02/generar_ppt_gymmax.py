"""
Generador de PowerPoint del Avance 02 de GymMax.
Imita fielmente el estilo visual del PPT del Avance 01:
- Fondo negro
- Color beige/tan (#d4a574) para títulos y acentos
- Marco redondeado beige en portada
- Títulos GRANDES en MAYÚSCULAS, italic, bold
- Decoración: 1 círculo arriba-izq + 3 círculos + flecha abajo-der
- Texto en italic con color claro
- Fuente Montserrat (con fallback Calibri)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ========= CONFIG =========
BASE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(BASE, "GymMax_Avance02.pptx")

# Paleta del Avance 01
BLACK = RGBColor(0x0a, 0x0a, 0x0a)
TAN = RGBColor(0xd4, 0xa5, 0x74)
TAN_BRIGHT = RGBColor(0xe8, 0xc4, 0x95)
WHITE = RGBColor(0xff, 0xff, 0xff)
WHITE_DIM = RGBColor(0xe0, 0xe0, 0xe0)
GRAY = RGBColor(0x9a, 0x9a, 0x9a)

TITLE_FONT = "Montserrat"
BODY_FONT = "Calibri"

# Tamaño widescreen 16:9
SLIDE_W = 13.333
SLIDE_H = 7.5


# ========= HELPERS =========

def set_bg(slide, color=BLACK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def no_line(shape):
    shape.line.fill.background()


def add_corner_decoration(slide, only_top=False):
    """1 círculo arriba-izq + 3 círculos y flecha abajo-der."""
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                              Inches(0.55), Inches(0.55),
                              Inches(0.32), Inches(0.32))
    c.fill.solid()
    c.fill.fore_color.rgb = TAN
    no_line(c)

    if only_top:
        return

    for i in range(3):
        c = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(12.65), Inches(5.5 + i * 0.55),
                                  Inches(0.28), Inches(0.28))
        c.fill.solid()
        c.fill.fore_color.rgb = TAN
        no_line(c)

    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                Inches(12.5), Inches(7.15),
                                Inches(0.55), Inches(0.18))
    arr.fill.solid()
    arr.fill.fore_color.rgb = TAN
    no_line(arr)


def add_footer(slide, page_num=None, total=None):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(7.15),
                                 Inches(7.0), Inches(0.25))
    tf = tb.text_frame
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "GYMMAX  ·  AVANCE 02  ·  DESARROLLO WEB INTEGRADO"
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = TAN
    run.font.name = TITLE_FONT

    if page_num and total:
        tb2 = slide.shapes.add_textbox(Inches(11.5), Inches(7.15),
                                      Inches(1.2), Inches(0.25))
        tf2 = tb2.text_frame
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.RIGHT
        run2 = p2.add_run()
        run2.text = f"{page_num} / {total}"
        run2.font.size = Pt(9)
        run2.font.color.rgb = GRAY
        run2.font.name = BODY_FONT


def add_title(slide, text, top=Inches(0.5), size=44, color=TAN, height=Inches(1.4)):
    tb = slide.shapes.add_textbox(Inches(0.8), top,
                                 Inches(SLIDE_W - 1.6), height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text.upper()
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.italic = True
    run.font.color.rgb = color
    run.font.name = TITLE_FONT
    return tb


def add_paragraph(slide, text, top, italic=True, color=WHITE_DIM, size=18,
                  bold=False, align=PP_ALIGN.LEFT, height=Inches(4.5)):
    tb = slide.shapes.add_textbox(Inches(0.9), top,
                                 Inches(SLIDE_W - 1.8), height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.italic = italic
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = BODY_FONT
    return tb


def add_bullets(slide, items, top, italic=True, color=WHITE_DIM, size=16,
                bullet="-", space_after=10, left=Inches(1.0), width=None):
    if width is None:
        width = Inches(SLIDE_W - 1.8)
    tb = slide.shapes.add_textbox(left, top, width, Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{bullet} {item}" if bullet else item
        run.font.size = Pt(size)
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = BODY_FONT
        p.space_after = Pt(space_after)
    return tb


def add_image_placeholder(slide, x, y, w, h, label="INSERTAR IMAGEN"):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0x22, 0x22, 0x22)
    rect.line.color.rgb = TAN
    rect.line.width = Pt(1.5)
    tf = rect.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "[ " + label + " ]"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.italic = True
    run.font.color.rgb = TAN
    run.font.name = BODY_FONT


def add_table(slide, x, y, w, h, headers, rows, header_size=11, body_size=10):
    cols = len(headers)
    rows_count = len(rows) + 1
    table_shape = slide.shapes.add_table(rows_count, cols, x, y, w, h)
    table = table_shape.table

    for i, hdr in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = TAN
        tf = cell.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = hdr
        run.font.size = Pt(header_size)
        run.font.bold = True
        run.font.color.rgb = BLACK
        run.font.name = BODY_FONT
        p.alignment = PP_ALIGN.CENTER

    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x1c, 0x1c, 0x1c)
            tf = cell.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(body_size)
            run.font.color.rgb = WHITE_DIM
            run.font.name = BODY_FONT


# ========= SLIDES =========

def slide_blank(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    set_bg(slide)
    return slide


def slide_portada(prs):
    slide = slide_blank(prs)

    margin_x = Inches(1.3)
    margin_y = Inches(0.7)
    border = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        margin_x, margin_y,
        Inches(SLIDE_W) - 2 * margin_x,
        Inches(SLIDE_H) - 2 * margin_y)
    border.fill.background()
    border.line.color.rgb = TAN
    border.line.width = Pt(2.5)

    c = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                              Inches(1.7), Inches(1.1),
                              Inches(0.3), Inches(0.3))
    c.fill.solid()
    c.fill.fore_color.rgb = TAN
    no_line(c)

    tb = slide.shapes.add_textbox(Inches(1.6), Inches(1.6),
                                 Inches(10), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "INFORME DE AVANCE 02"
    run.font.size = Pt(46)
    run.font.bold = True
    run.font.italic = True
    run.font.color.rgb = TAN
    run.font.name = TITLE_FONT

    tb2 = slide.shapes.add_textbox(Inches(1.6), Inches(2.6),
                                  Inches(10), Inches(1.2))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = "GYMMAX"
    run2.font.size = Pt(72)
    run2.font.bold = True
    run2.font.italic = True
    run2.font.color.rgb = TAN_BRIGHT
    run2.font.name = TITLE_FONT

    tb3 = slide.shapes.add_textbox(Inches(1.6), Inches(3.85),
                                  Inches(10), Inches(0.6))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.LEFT
    run3 = p3.add_run()
    run3.text = "Sistema Integral de Gestión de Gimnasios"
    run3.font.size = Pt(20)
    run3.font.italic = True
    run3.font.color.rgb = WHITE_DIM
    run3.font.name = BODY_FONT

    tb4 = slide.shapes.add_textbox(Inches(1.6), Inches(4.7),
                                  Inches(10), Inches(0.5))
    tf4 = tb4.text_frame
    p4 = tf4.paragraphs[0]
    run4 = p4.add_run()
    run4.text = "INTEGRANTES:"
    run4.font.size = Pt(16)
    run4.font.bold = True
    run4.font.italic = True
    run4.font.color.rgb = TAN
    run4.font.name = TITLE_FONT

    integrantes = [
        "CHOQUE ANCHANTE NIURKA YASBETH",
        "DIAZ CULQUI NEHEMIAS",
        "PURIZACA IPANAQUE DENNYS MARLON",
        "SORIA CHAVEZ IAN",
        "TORRE ESCOBAR OLIVER",
        "VALLADOLID LLENQUE ALEXANDER",
    ]

    tb5 = slide.shapes.add_textbox(Inches(1.6), Inches(5.15),
                                  Inches(9), Inches(1.8))
    tf5 = tb5.text_frame
    tf5.word_wrap = True
    for i, nombre in enumerate(integrantes):
        if i == 0:
            p = tf5.paragraphs[0]
        else:
            p = tf5.add_paragraph()
        run = p.add_run()
        run.text = "- " + nombre
        run.font.size = Pt(13)
        run.font.bold = True
        run.font.italic = True
        run.font.color.rgb = TAN_BRIGHT
        run.font.name = TITLE_FONT
        p.space_after = Pt(2)

    logo = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(11.2), Inches(5.7),
                                 Inches(0.8), Inches(0.8))
    logo.fill.solid()
    logo.fill.fore_color.rgb = RGBColor(0xc8, 0x10, 0x2e)
    no_line(logo)
    tf = logo.text_frame
    tf.margin_top = Inches(0.18)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "UTP"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = WHITE

    tbF = slide.shapes.add_textbox(Inches(1.6), Inches(6.6),
                                  Inches(8), Inches(0.4))
    tfF = tbF.text_frame
    pF = tfF.paragraphs[0]
    runF = pF.add_run()
    runF.text = "DESARROLLO WEB INTEGRADO  ·  UTP  ·  MAYO 2026"
    runF.font.size = Pt(11)
    runF.font.bold = True
    runF.font.italic = True
    runF.font.color.rgb = TAN
    runF.font.name = TITLE_FONT


def slide_integrantes(prs):
    slide = slide_blank(prs)
    add_corner_decoration(slide)
    add_footer(slide, 2, 26)

    border = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(1.4), Inches(0.9),
                                   Inches(10.5), Inches(5.7))
    border.fill.background()
    border.line.color.rgb = TAN
    border.line.width = Pt(2)

    add_title(slide, "INFORME DE AVANCE 02 — GYMMAX",
              top=Inches(1.3), size=32)

    tb = slide.shapes.add_textbox(Inches(1.9), Inches(2.5),
                                 Inches(9), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "INTEGRANTES:"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.italic = True
    run.font.color.rgb = TAN
    run.font.name = TITLE_FONT

    integrantes = [
        "CHOQUE ANCHANTE NIURKA YASBETH",
        "DIAZ CULQUI NEHEMIAS",
        "PURIZACA IPANAQUE DENNYS MARLON",
        "SORIA CHAVEZ IAN",
        "TORRE ESCOBAR OLIVER",
        "VALLADOLID LLENQUE ALEXANDER",
    ]
    add_bullets(slide, integrantes, top=Inches(3.1),
                size=18, color=TAN_BRIGHT, italic=True,
                left=Inches(2.2))


def slide_agenda(prs):
    slide = slide_blank(prs)
    add_corner_decoration(slide)
    add_footer(slide, 3, 26)
    add_title(slide, "Agenda", size=54)

    agenda = [
        "Introducción y problemática",
        "Objetivos y alcance",
        "Características del producto",
        "Product Backlog y User Stories",
        "Wireframes y flujos de navegación",
        "Modelo de Base de Datos (Lógico y Físico)",
        "Diagrama de Clases UML",
        "Tecnologías utilizadas",
        "Próximos pasos y conclusiones",
    ]
    add_bullets(slide, [f"{i+1}. {t}" for i, t in enumerate(agenda)],
                top=Inches(2.1), size=20, italic=True,
                color=TAN_BRIGHT, bullet="", space_after=14,
                left=Inches(1.5))


def slide_producto(prs):
    slide = slide_blank(prs)
    add_corner_decoration(slide)
    add_footer(slide, 4, 26)
    add_title(slide, "El Producto", size=54)

    add_paragraph(slide,
        "GymMax es una plataforma web integral construida bajo la arquitectura "
        "Java Enterprise Edition (JEE) que centraliza la gestión de cadenas "
        "de gimnasios: sedes, socios, membresías digitales y reservas de clases.",
        top=Inches(2.0), italic=True, size=17, height=Inches(2.0))

    add_bullets(slide, [
        "Plataforma centralizada multi-sede",
        "Construida sobre Java EE + MySQL + Bootstrap 5.3",
        "Diseño Mobile First — accesible desde cualquier dispositivo",
        "Gestión completa del ciclo de vida del socio",
    ], top=Inches(4.2), size=16, italic=True, color=TAN_BRIGHT)


def slide_problematica(prs):
    slide = slide_blank(prs)
    add_corner_decoration(slide)
    add_footer(slide, 5, 26)
    add_title(slide, "Introducción y Problemática", size=40)

    add_paragraph(slide,
        "Los gimnasios en Perú enfrentan serios desafíos en la gestión manual de "
        "membresías, control de accesos y administración de sedes. Esta falta de "
        "automatización genera pérdida de ingresos, sobrecarga administrativa y "
        "mala experiencia para el socio.",
        top=Inches(2.0), italic=True, size=15, height=Inches(1.8))

    add_paragraph(slide, "Problemas identificados:",
                  top=Inches(4.0), italic=True, bold=True,
                  color=TAN, size=16, height=Inches(0.4))

    add_bullets(slide, [
        "Pérdida de ingresos por nulo control de vencimientos",
        "Sobrecarga administrativa y procesos manuales",
        "Procesos lentos y largas colas en recepción",
        "Mala experiencia del socio (no puede gestionar desde celular)",
        "Sistemas existentes con lentitud y caídas en horas pico",
    ], top=Inches(4.5), size=15, italic=True, color=TAN_BRIGHT,
       space_after=8, left=Inches(1.2))


def slide_objetivo_gral(prs):
    slide = slide_blank(prs)
    add_corner_decoration(slide)
    add_footer(slide, 6, 26)
    add_title(slide, "Objetivo del Proyecto General", size=40)

    add_paragraph(slide,
        "Desarrollar una aplicación web integral bajo la arquitectura Java "
        "Enterprise Edition (JEE) que centralice la visualización de sedes, el "
        "registro de socios y la venta de membresías digitales. Para lograrlo "
        "se implementarán módulos de consulta de sedes con horarios y servicios, "
        "un formulario de registro responsive con validación en tiempo real "
        "almacenado en MySQL, y un sistema de gestión de planes que permita "
        "renovaciones autónomas. Todo el diseño se basará en una interfaz "
        "\"Mobile First\" utilizando Bootstrap 5.3 para garantizar accesibilidad total.",
        top=Inches(2.0), italic=True, bold=True, size=16,
        color=TAN_BRIGHT, height=Inches(5.0))


def slide_objetivos_esp(prs):
    slide = slide_blank(prs)
    add_corner_decoration(slide)
    add_footer(slide, 7, 26)
    add_title(slide, "Objetivos Específicos", size=44)

    objs = [
        "Implementar un módulo que permita visualizar todas las sedes disponibles, "
        "incluyendo ubicación, horarios, servicios y cupos por clase.",
        "Desarrollar un formulario web intuitivo y responsive para registro de "
        "socios con validación en tiempo real y almacenamiento seguro en MySQL.",
        "Crear un sistema de gestión de planes (básico, premium, anual) con "
        "renovación autónoma, cálculo de vencimiento y comprobante digital.",
        "Diseñar una interfaz Mobile First con Bootstrap 5.3, intuitiva y "
        "accesible desde cualquier dispositivo.",
        "Implementar panel administrativo con KPIs en tiempo real, CRUD de "
        "entidades y reportes exportables a Excel y PDF.",
    ]
    add_bullets(slide, objs, top=Inches(1.8), size=14, italic=True,
                color=TAN_BRIGHT, space_after=14, left=Inches(1.0))


def slide_alcance(prs):
    slide = slide_blank(prs)
    add_corner_decoration(slide)
    add_footer(slide, 8, 26)
    add_title(slide, "Alcance del Proyecto", size=44)

    add_paragraph(slide,
        "Plataforma web centralizada que permite la gestión completa del ciclo "
        "de vida del socio: desde el registro inicial y la compra de membresías "
        "hasta la reserva de clases por sede. Incluye panel administrativo, "
        "reportes en tiempo real y arquitectura JEE escalable.",
        top=Inches(1.9), italic=True, size=14, height=Inches(2.2))

    add_paragraph(slide, "SÍ incluye:", top=Inches(4.2),
                  italic=True, bold=True, color=TAN, size=14, height=Inches(0.4))
    add_bullets(slide, [
        "Multi-sede, multi-rol (Socio + Admin)",
        "Membresías digitales y pagos online",
        "Dashboard con KPIs y reportes exportables",
    ], top=Inches(4.6), size=12, italic=True, color=TAN_BRIGHT,
       space_after=4, left=Inches(1.0))

    add_paragraph(slide, "NO incluye:", top=Inches(5.95),
                  italic=True, bold=True, color=TAN, size=14, height=Inches(0.4))
    add_bullets(slide, [
        "Pagos en efectivo (solo digitales)",
        "App móvil nativa (Android/iOS)",
        "Integración biométrica de accesos",
    ], top=Inches(6.35), size=12, italic=True, color=TAN_BRIGHT,
       space_after=2, left=Inches(1.0))


def slide_caracteristicas(prs):
    slide = slide_blank(prs)
    add_corner_decoration(slide)
    add_footer(slide, 9, 26)
    add_title(slide, "Características del Producto", size=38)

    add_bullets(slide, [
        "Gestión multi-sede centralizada",
        "Sistema de roles diferenciados (Socio / Administrador)",
        "Registro y autenticación con validación en tiempo real",
        "Membresías digitales (Básico, Premium, Anual)",
        "Reserva de clases online con control de cupos",
        "Pagos digitales: Yape, Plin, Tarjeta",
        "Dashboard administrativo con KPIs en tiempo real",
        "Reportes filtrables y exportación a Excel / PDF",
        "Diseño Mobile First con Bootstrap 5.3",
        "Arquitectura por capas MVC + DAO + Facade",
    ], top=Inches(1.95), size=14, italic=True, color=TAN_BRIGHT,
       space_after=6)


def slide_backlog(prs):
    slide = slide_blank(prs)
    add_corner_decoration(slide)
    add_footer(slide, 10, 26)
    add_title(slide, "Product Backlog", size=44)

    add_paragraph(slide,
        "28 requerimientos funcionales + 12 no funcionales distribuidos en 5 sprints.",
        top=Inches(1.8), size=14, italic=True, height=Inches(0.5))

    add_table(slide, Inches(1.5), Inches(2.5),
              Inches(10.3), Inches(3.5),
              ["Sprint", "Foco", "Items", "SP"],
              [
                  ("1", "Autenticación + base socio", "RF-01, 02, 04, 05, 06, 18", "22"),
                  ("2", "Sedes, clases y reservas", "RF-07, 08, 09, 10, 11, 19, 20, 21", "37"),
                  ("3", "Membresías y pagos", "RF-03, 12, 13, 14, 15, 16, 17, 22", "41"),
                  ("4", "Reportes y asistencias", "RF-23, 24, 25, 26, 28", "26"),
                  ("5", "Notificaciones", "RF-27", "8"),
              ],
              header_size=12, body_size=11)

    add_paragraph(slide, "Total: 134 Story Points en 5 sprints.",
                  top=Inches(6.3), size=12, italic=True, color=TAN)


def slide_user_stories(prs):
    slide = slide_blank(prs)
    add_corner_decoration(slide)
    add_footer(slide, 11, 26)
    add_title(slide, "User Stories Destacadas", size=40)

    add_paragraph(slide, "US-08 — Reservar clase",
                  top=Inches(1.9), bold=True, italic=True,
                  color=TAN, size=14, height=Inches(0.4))
    add_paragraph(slide,
        "Como socio con membresía activa quiero reservar una clase para asegurar "
        "mi cupo. Verifica membresía + cupos + duplicado.",
        top=Inches(2.35), italic=True, size=12,
        color=WHITE_DIM, height=Inches(0.8))

    add_paragraph(slide, "US-14 — Dashboard administrativo",
                  top=Inches(3.3), bold=True, italic=True,
                  color=TAN, size=14, height=Inches(0.4))
    add_paragraph(slide,
        "Como admin quiero ver KPIs en tiempo real para decidir. 4 KPIs + "
        "variación porcentual + gráfico de ingresos últimos 6 meses.",
        top=Inches(3.75), italic=True, size=12,
        color=WHITE_DIM, height=Inches(0.8))

    add_paragraph(slide, "US-17 — Generación de reportes",
                  top=Inches(4.7), bold=True, italic=True,
                  color=TAN, size=14, height=Inches(0.4))
    add_paragraph(slide,
        "Como admin quiero reportes filtrables (4 tipos) con exportación a "
        "Excel y PDF para compartirlos con la gerencia.",
        top=Inches(5.15), italic=True, size=12,
        color=WHITE_DIM, height=Inches(0.8))

    add_paragraph(slide,
        "El informe Word contiene las 18 historias completas con criterios de aceptación.",
        top=Inches(6.4), italic=True, size=11, color=GRAY)


def slide_flujo_socio(prs):
    slide = slide_blank(prs)
    add_corner_decoration(slide)
    add_footer(slide, 12, 26)
    add_title(slide, "Flujo de Navegación — Socio", size=36)

    add_image_placeholder(slide, Inches(1.6), Inches(2.0),
                         Inches(10), Inches(4.0),
                         "INSERTAR DIAGRAMA DE FLUJO SOCIO (del Avance 01)")
    add_paragraph(slide,
        "Login → validar credenciales → Dashboard del socio → "
        "Ver sedes / Reservar clase / Mis reservas",
        top=Inches(6.3), italic=True, size=12, color=TAN_BRIGHT)


def slide_flujo_admin(prs):
    slide = slide_blank(prs)
    add_corner_decoration(slide)
    add_footer(slide, 13, 26)
    add_title(slide, "Flujo de Navegación — Administrador", size=32)

    add_image_placeholder(slide, Inches(1.6), Inches(2.0),
                         Inches(10), Inches(4.0),
                         "INSERTAR DIAGRAMA DE FLUJO ADMIN (del Avance 01)")
    add_paragraph(slide,
        "Login → validar rol ADMIN → Dashboard administrativo → "
        "Gestión socios / Gestión sedes / Reportes",
        top=Inches(6.3), italic=True, size=12, color=TAN_BRIGHT)


def slide_mockups_mobile_1(prs):
    slide = slide_blank(prs)
    add_corner_decoration(slide)
    add_footer(slide, 14, 26)
    add_title(slide, "Wireframes Mobile (1/2)", size=36)

    add_image_placeholder(slide, Inches(0.9), Inches(2.0),
                         Inches(3.7), Inches(4.5), "Login")
    add_image_placeholder(slide, Inches(4.85), Inches(2.0),
                         Inches(3.7), Inches(4.5), "Registro")
    add_image_placeholder(slide, Inches(8.8), Inches(2.0),
                         Inches(3.7), Inches(4.5), "Dashboard Socio")


def slide_mockups_mobile_2(prs):
    slide = slide_blank(prs)
    add_corner_decoration(slide)
    add_footer(slide, 15, 26)
    add_title(slide, "Wireframes Mobile (2/2)", size=36)

    add_image_placeholder(slide, Inches(0.9), Inches(2.0),
                         Inches(3.7), Inches(4.5), "Listado Sedes")
    add_image_placeholder(slide, Inches(4.85), Inches(2.0),
                         Inches(3.7), Inches(4.5), "Reservar Clase")
    add_image_placeholder(slide, Inches(8.8), Inches(2.0),
                         Inches(3.7), Inches(4.5), "Mis Reservas")


def slide_mockups_desktop(prs):
    slide = slide_blank(prs)
    add_corner_decoration(slide)
    add_footer(slide, 16, 26)
    add_title(slide, "Wireframes Desktop", size=40)

    add_image_placeholder(slide, Inches(0.9), Inches(2.0),
                         Inches(5.9), Inches(4.5),
                         "Dashboard Admin (KPIs)")
    add_image_placeholder(slide, Inches(6.95), Inches(2.0),
                         Inches(5.6), Inches(4.5),
                         "Gestión de Socios")


def slide_reportes_desktop(prs):
    slide = slide_blank(prs)
    add_corner_decoration(slide)
    add_footer(slide, 17, 26)
    add_title(slide, "Reportes — Versión Desktop", size=38)

    add_image_placeholder(slide, Inches(2.0), Inches(2.0),
                         Inches(9.3), Inches(4.5),
                         "Generación de Reportes con gráfico y tabla")


def slide_bd_logico(prs):
    slide = slide_blank(prs)
    add_corner_decoration(slide)
    add_footer(slide, 18, 26)
    add_title(slide, "Modelo de Base de Datos — Lógico", size=34)

    add_image_placeholder(slide, Inches(1.0), Inches(1.9),
                         Inches(11.3), Inches(4.0),
                         "Diagrama Lógico (Draw.io / Lucid Chart)")
    add_paragraph(slide,
        "10 entidades · 10 relaciones · sin tipos de dato · orientado al negocio",
        top=Inches(6.2), italic=True, size=12, color=TAN_BRIGHT)


def slide_bd_fisico(prs):
    slide = slide_blank(prs)
    add_corner_decoration(slide)
    add_footer(slide, 19, 26)
    add_title(slide, "Modelo de Base de Datos — Físico", size=34)

    add_image_placeholder(slide, Inches(1.0), Inches(1.9),
                         Inches(11.3), Inches(4.0),
                         "Diagrama Físico MySQL 8.0 (Draw.io / Lucid Chart)")
    add_paragraph(slide,
        "Motor InnoDB · charset utf8mb4 · 10 tablas con tipos MySQL · "
        "PKs, FKs, ENUMs, índices",
        top=Inches(6.2), italic=True, size=12, color=TAN_BRIGHT)


def slide_tecnologias(prs):
    slide = slide_blank(prs)
    add_corner_decoration(slide)
    add_footer(slide, 20, 26)
    add_title(slide, "Stack Tecnológico", size=44)

    add_table(slide, Inches(2.0), Inches(2.0),
              Inches(9.3), Inches(4.7),
              ["Capa", "Tecnología"],
              [
                  ("Lenguaje", "Java 11 + Jakarta EE 10"),
                  ("Frontend", "JSP + HTML5 + Bootstrap 5.3"),
                  ("Controlador", "Servlets Jakarta 6.0"),
                  ("Lógica negocio", "Facades (POJO)"),
                  ("Persistencia", "JDBC + DAO + Connector/J"),
                  ("Base de datos", "MySQL 8.0 Community"),
                  ("Servidor", "Apache Tomcat 10.1"),
                  ("IDE", "NetBeans 21+"),
                  ("Modelado", "Draw.io / Lucid Chart"),
                  ("Metodología", "Scrum (sprints 2 semanas)"),
              ],
              header_size=14, body_size=12)


def slide_uml_1(prs):
    slide = slide_blank(prs)
    add_corner_decoration(slide)
    add_footer(slide, 21, 26)
    add_title(slide, "Diagrama de Clases UML (1/2)", size=36)

    add_image_placeholder(slide, Inches(1.0), Inches(1.9),
                         Inches(11.3), Inches(4.0),
                         "DTOs + DAOs (interfaces + implementaciones)")
    add_paragraph(slide,
        "10 DTOs (modelo de datos) + 9 interfaces DAO + 9 implementaciones JDBC. "
        "Programación a interfaces para desacoplar persistencia.",
        top=Inches(6.2), italic=True, size=12, color=TAN_BRIGHT,
        height=Inches(0.8))


def slide_uml_2(prs):
    slide = slide_blank(prs)
    add_corner_decoration(slide)
    add_footer(slide, 22, 26)
    add_title(slide, "Diagrama de Clases UML (2/2)", size=36)

    add_image_placeholder(slide, Inches(1.0), Inches(1.9),
                         Inches(11.3), Inches(4.0),
                         "Facades + Servlets + Conexion")
    add_paragraph(slide,
        "5 Facades (lógica de negocio) · 5 Servlets (controladores MVC) · "
        "1 Conexion singleton. Total: 39 clases.",
        top=Inches(6.2), italic=True, size=12, color=TAN_BRIGHT,
        height=Inches(0.8))


def slide_diccionario(prs):
    slide = slide_blank(prs)
    add_corner_decoration(slide)
    add_footer(slide, 23, 26)
    add_title(slide, "Diccionario de Clases", size=44)

    add_table(slide, Inches(3.0), Inches(2.0),
              Inches(7.3), Inches(4.2),
              ["Capa", "Cantidad"],
              [
                  ("DTO", "10 clases"),
                  ("DAO interfaces", "9"),
                  ("DAO implementaciones", "9"),
                  ("Facades", "5"),
                  ("Controllers (Servlets)", "5"),
                  ("Utilitaria (Conexion)", "1"),
                  ("TOTAL", "39 clases"),
              ],
              header_size=14, body_size=12)

    add_paragraph(slide,
        "El informe Word contiene la descripción detallada de cada clase con "
        "atributos, métodos, parámetros y reglas de negocio.",
        top=Inches(6.4), italic=True, size=12, color=TAN_BRIGHT,
        align=PP_ALIGN.CENTER)


def slide_patrones(prs):
    slide = slide_blank(prs)
    add_corner_decoration(slide)
    add_footer(slide, 24, 26)
    add_title(slide, "Patrones de Diseño Aplicados", size=38)

    add_bullets(slide, [
        "MVC  ·  Model (DTO) — View (JSP) — Controller (Servlet)",
        "DAO  ·  Aislamiento del acceso a datos mediante interfaces + implementación",
        "Facade  ·  Encapsulamiento de la lógica de negocio entre Controller y DAO",
        "Singleton  ·  Helper único de conexión JDBC (Conexion.java)",
        "DTO  ·  Objetos planos que transportan datos entre capas",
        "Programación a interfaces  ·  Bajo acoplamiento, fácil testing y mantenimiento",
    ], top=Inches(2.1), size=16, italic=True, color=TAN_BRIGHT,
       space_after=14, left=Inches(1.0))


def slide_proximos(prs):
    slide = slide_blank(prs)
    add_corner_decoration(slide)
    add_footer(slide, 25, 26)
    add_title(slide, "Hoja de Ruta — Entrega Final", size=38)

    add_bullets(slide, [
        "Sprint 1-2  (Semanas 11-12)  ·  Autenticación + dashboard socio + listado de sedes",
        "Sprint 3    (Semana 13)        ·  Reservas de clases + membresías y pagos",
        "Sprint 4    (Semana 14)        ·  Panel administrativo + reportes + exportación",
        "Sprint 5    (Semana 15)        ·  Notificaciones + pruebas integrales + documentación final",
    ], top=Inches(2.5), size=15, italic=True, color=TAN_BRIGHT,
       space_after=20, left=Inches(0.8))


def slide_conclusiones(prs):
    slide = slide_blank(prs)
    add_corner_decoration(slide)
    add_footer(slide, 26, 26)
    add_title(slide, "Conclusiones", size=54)

    add_bullets(slide, [
        "Diseño completo del sistema cubierto: 18 user stories, 40 requerimientos",
        "Modelo de datos validado: 10 tablas normalizadas + script SQL listo",
        "Arquitectura por capas definida: 39 clases organizadas en 5 paquetes",
        "Stack tecnológico consolidado: Java EE + MySQL + Bootstrap + Tomcat",
        "Equipo alineado y listo para iniciar Sprint 1 a partir de la Semana 11",
    ], top=Inches(2.3), size=16, italic=True, color=TAN_BRIGHT,
       space_after=18, left=Inches(0.9))


def slide_gracias(prs):
    slide = slide_blank(prs)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.7),
                                 Inches(12.3), Inches(2.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "GRACIAS"
    run.font.size = Pt(140)
    run.font.bold = True
    run.font.italic = True
    run.font.color.rgb = TAN
    run.font.name = TITLE_FONT


# ========= MAIN =========

def main():
    print("Generando PowerPoint estilo Avance 01...")
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    slide_portada(prs)
    slide_integrantes(prs)
    slide_agenda(prs)
    slide_producto(prs)
    slide_problematica(prs)
    slide_objetivo_gral(prs)
    slide_objetivos_esp(prs)
    slide_alcance(prs)
    slide_caracteristicas(prs)
    slide_backlog(prs)
    slide_user_stories(prs)
    slide_flujo_socio(prs)
    slide_flujo_admin(prs)
    slide_mockups_mobile_1(prs)
    slide_mockups_mobile_2(prs)
    slide_mockups_desktop(prs)
    slide_reportes_desktop(prs)
    slide_bd_logico(prs)
    slide_bd_fisico(prs)
    slide_tecnologias(prs)
    slide_uml_1(prs)
    slide_uml_2(prs)
    slide_diccionario(prs)
    slide_patrones(prs)
    slide_proximos(prs)
    slide_conclusiones(prs)
    slide_gracias(prs)

    prs.save(OUT)
    print(f"  -> {OUT}")
    print(f"\nTotal: {len(prs.slides)} slides\n")
    print("=== PowerPoint generado ===")


if __name__ == "__main__":
    main()
